"""Update data-driven tables on the DM review pptx.

Currently handles:
- Slide 2 / "Tabel 13": the SERP rankings table. Columns are
  ``Type URL | {prev_month} | {last_month} | Delta`` and rows are the
  four URL types Cat-url / C-url / PLP / R-url. Data is sourced from the
  workbook's `serp` tab (wide layout with one column per month).
- Slide 2 / "Tabel 25" (Visits) and "Tabel 27" (Revenue): the target /
  behaald cards. Targets come from `seo_targets.xlsx` sheet ``2026``,
  achieved from the SEO row of the latest month in `visits_omzet`.

The pptx is matched by inspecting cells: any 4-row × 4-col table on slide 2
whose first column reads ``Type URL`` is the SERP table. The target cards
are 3-column ``Kanaal | Target | Behaald`` tables — the one whose Target
column contains a € sign is the revenue card; the other is visits.
"""
from __future__ import annotations

import datetime as _dt
import logging
import os
from copy import copy
from typing import Dict, List, Optional, Tuple

from openpyxl import load_workbook
from pptx import Presentation

logger = logging.getLogger(__name__)

TARGETS_EXCEL_PATH = "/mnt/c/Users/JoepvanSchagen/OneDrive - Beslist.nl BV/seo_targets.xlsx"
TARGETS_SHEET = "2026"
TARGETS_VISITS_ROW = 8
TARGETS_REVENUE_ROW = 6
TARGETS_MONTH_COL_BASE = 2  # Jan = col 3, so base + month_number

DUTCH_MONTH_FULL = {
    1: "Januari", 2: "Februari", 3: "Maart", 4: "April", 5: "Mei", 6: "Juni",
    7: "Juli", 8: "Augustus", 9: "September", 10: "Oktober", 11: "November", 12: "December",
}

URL_TYPE_ORDER = ("Cat-url", "C-url", "PLP", "R-url")


def _check_pptx_lock(path: str) -> Optional[str]:
    d, base = os.path.split(path)
    lock = os.path.join(d, "~$" + base)
    if os.path.exists(lock):
        return (f"The presentation is currently open. Please close "
                f"'{base}' in PowerPoint and try again.")
    if not os.path.exists(path):
        return f"PowerPoint file not found at {path}"
    return None


def _read_serp_last_two_months(excel_path: str, target_yyyymm: Optional[int] = None
                               ) -> Optional[Tuple[str, str, Dict[str, float], Dict[str, float]]]:
    """Return (prev_month_name, last_month_name, prev_values_by_type, last_values_by_type).

    Reads two adjacent month columns from the `serp` tab (those whose header is a
    Dutch month name), skipping the Delta column. When `target_yyyymm` is given,
    the "last" month is that month's column (the rightmost one with that name, to
    handle the annually-recycling header) and "prev" is the month column directly
    to its left — so a May review reports April | Mei even if a later Juni column
    already exists. Without it, the two rightmost columns are used.
    """
    wb = load_workbook(excel_path, data_only=True, read_only=True)
    if "serp" not in wb.sheetnames:
        return None
    ws = wb["serp"]
    valid_months = set(DUTCH_MONTH_FULL.values())
    month_cols: List[Tuple[int, str]] = []
    for c in range(2, ws.max_column + 1):
        v = ws.cell(row=1, column=c).value
        if isinstance(v, str) and v.strip().capitalize() in valid_months:
            month_cols.append((c, v.strip().capitalize()))
    if len(month_cols) < 2:
        return None

    last_idx = len(month_cols) - 1  # default: rightmost
    if target_yyyymm is not None:
        target_name = DUTCH_MONTH_FULL[target_yyyymm % 100]
        # rightmost column matching the target month (names recycle each year)
        match = [i for i, (_, name) in enumerate(month_cols) if name == target_name]
        if match and match[-1] >= 1:
            last_idx = match[-1]
    prev_col, prev_name = month_cols[last_idx - 1]
    last_col, last_name = month_cols[last_idx]
    prev_vals: Dict[str, float] = {}
    last_vals: Dict[str, float] = {}
    for r in range(2, ws.max_row + 1):
        label = ws.cell(row=r, column=1).value
        if not isinstance(label, str):
            continue
        label = label.strip()
        if label not in URL_TYPE_ORDER:
            continue
        p = ws.cell(row=r, column=prev_col).value
        l = ws.cell(row=r, column=last_col).value
        if isinstance(p, (int, float)):
            prev_vals[label] = float(p)
        if isinstance(l, (int, float)):
            last_vals[label] = float(l)
    return prev_name, last_name, prev_vals, last_vals


def _fmt_position(v: Optional[float]) -> str:
    if v is None:
        return ""
    # Dutch comma decimal, 2 decimals.
    return f"{v:.2f}".replace(".", ",")


def _fmt_delta(prev: Optional[float], last: Optional[float]) -> str:
    if not prev or last is None:
        return ""
    pct = (last - prev) / prev * 100
    # Round to integer percent, keep negative sign, no plus sign for positives
    # (matches the convention already used in the existing table).
    return f"{round(pct):.0f}%"


def _set_cell_text(cell, text: str) -> None:
    """Replace cell text while preserving font formatting from the first run."""
    tf = cell.text_frame
    if not tf.paragraphs:
        cell.text = text
        return
    para = tf.paragraphs[0]
    if not para.runs:
        cell.text = text
        return
    first_run = para.runs[0]
    # Capture font properties
    src_font = first_run.font
    # Wipe paragraph but keep the run
    for run in para.runs[1:]:
        run.text = ""
    first_run.text = text
    # Restore font (text reassignment usually preserves run-level font)
    if src_font.size is not None:
        first_run.font.size = src_font.size
    if src_font.bold is not None:
        first_run.font.bold = src_font.bold


def _find_serp_table(slide):
    """Return the SERP rankings table on the slide (4 cols, first cell 'Type URL')."""
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        if len(tbl.columns) != 4 or len(tbl.rows) < 2:
            continue
        first_cell = tbl.cell(0, 0).text.strip().lower()
        if first_cell == "type url":
            return shape, tbl
    return None, None


def _read_seo_latest_month_actuals(excel_path: str) -> Optional[Tuple[_dt.date, int, float]]:
    """Return (month_date, seo_visits, seo_omzet) for the most recent month in visits_omzet."""
    wb = load_workbook(excel_path, data_only=True, read_only=True)
    if "visits_omzet" not in wb.sheetnames:
        return None
    ws = wb["visits_omzet"]
    best: Optional[Tuple[_dt.date, int, float]] = None
    for r in range(2, ws.max_row + 1):
        d = ws.cell(row=r, column=1).value
        v = ws.cell(row=r, column=2).value
        o = ws.cell(row=r, column=3).value
        k = ws.cell(row=r, column=4).value
        if not (isinstance(d, _dt.datetime) and k == "SEO"):
            continue
        if v is None:
            continue
        date_only = d.date()
        if best is None or date_only > best[0]:
            best = (date_only, int(v), float(o or 0))
    return best


def _read_seo_month_actuals(excel_path: str, yyyymm: int) -> Optional[Tuple[_dt.date, int, float]]:
    """Return (month_date, seo_visits, seo_omzet) for a specific month in visits_omzet.

    Used when slide 2 is built for a month other than the most recent one, so the
    target/behaald cards read that month's row instead of the latest. Returns None
    if there's no SEO row for that month yet."""
    wb = load_workbook(excel_path, data_only=True, read_only=True)
    if "visits_omzet" not in wb.sheetnames:
        return None
    ws = wb["visits_omzet"]
    y, m = divmod(yyyymm, 100)
    for r in range(2, ws.max_row + 1):
        d = ws.cell(row=r, column=1).value
        v = ws.cell(row=r, column=2).value
        o = ws.cell(row=r, column=3).value
        k = ws.cell(row=r, column=4).value
        if not (isinstance(d, _dt.datetime) and k == "SEO"):
            continue
        if v is None:
            continue
        if d.year == y and d.month == m:
            return (d.date(), int(v), float(o or 0))
    return None


def _read_targets_for_month(month: int) -> Optional[Tuple[float, float]]:
    """Return (visits_target, omzet_target) for the given month (1..12) from seo_targets.xlsx."""
    if not os.path.exists(TARGETS_EXCEL_PATH):
        return None
    wb = load_workbook(TARGETS_EXCEL_PATH, data_only=True, read_only=True)
    if TARGETS_SHEET not in wb.sheetnames:
        return None
    ws = wb[TARGETS_SHEET]
    col = TARGETS_MONTH_COL_BASE + month  # Jan(1) → col 3
    visits = ws.cell(row=TARGETS_VISITS_ROW, column=col).value
    omzet = ws.cell(row=TARGETS_REVENUE_ROW, column=col).value
    if visits is None or omzet is None:
        return None
    return float(visits), float(omzet)


def _fmt_int_dutch(v: float) -> str:
    """Thousands separator with dot, no decimals: 2.881.649"""
    return f"{int(round(v)):,}".replace(",", ".")


def _fmt_money_dutch(v: float) -> str:
    return f"€ {_fmt_int_dutch(v)}"


def _fmt_pct(v: float) -> str:
    return f"{round(v):.0f}%"


def _find_target_cards(slide) -> Tuple[Optional[object], Optional[object]]:
    """Return (visits_card, revenue_card). Revenue card's row-1 col-1 contains €."""
    visits = None
    revenue = None
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        if len(tbl.columns) != 3 or len(tbl.rows) != 2:
            continue
        header = [tbl.cell(0, c).text.strip().lower() for c in range(3)]
        if header != ["kanaal", "target", "behaald"]:
            continue
        target_cell = tbl.cell(1, 1).text
        if "€" in target_cell:
            revenue = shape
        else:
            visits = shape
    return visits, revenue


def update_target_cards(pptx_path: str, excel_path: str, slide_index: int = 1,
                        target_yyyymm: Optional[int] = None) -> Dict:
    """Update slide-2 Visits + Revenue target/behaald cards.

    When `target_yyyymm` is given, the achieved values are read from that month's
    SEO row in visits_omzet (falling back to the latest month if that row doesn't
    exist yet). Without it, the latest month is used."""
    lock_err = _check_pptx_lock(pptx_path)
    if lock_err:
        return {"status": "error", "error": lock_err}

    actuals = None
    if target_yyyymm is not None:
        actuals = _read_seo_month_actuals(excel_path, target_yyyymm)
    if actuals is None:
        actuals = _read_seo_latest_month_actuals(excel_path)
    if actuals is None:
        return {"status": "error", "error": "No SEO data in visits_omzet"}
    month_date, seo_visits, seo_omzet = actuals
    targets = _read_targets_for_month(month_date.month)
    if targets is None:
        return {"status": "error",
                "error": f"No targets for month {month_date.month} in {TARGETS_EXCEL_PATH}"}
    visits_target, omzet_target = targets

    prs = Presentation(pptx_path)
    if slide_index >= len(prs.slides):
        return {"status": "error",
                "error": f"Slide index {slide_index} out of range"}
    slide = prs.slides[slide_index]
    visits_shape, revenue_shape = _find_target_cards(slide)
    if visits_shape is None or revenue_shape is None:
        return {"status": "error", "error": "Could not find both Visits and Revenue target cards"}

    visits_pct = (seo_visits / visits_target * 100) if visits_target else 0
    omzet_pct = (seo_omzet / omzet_target * 100) if omzet_target else 0

    # Visits card
    vtbl = visits_shape.table
    _set_cell_text(vtbl.cell(1, 1), _fmt_int_dutch(visits_target))
    _set_cell_text(vtbl.cell(1, 2), _fmt_pct(visits_pct))
    # Revenue card
    rtbl = revenue_shape.table
    _set_cell_text(rtbl.cell(1, 1), _fmt_money_dutch(omzet_target))
    _set_cell_text(rtbl.cell(1, 2), _fmt_pct(omzet_pct))

    lock_err = _check_pptx_lock(pptx_path)
    if lock_err:
        return {"status": "error", "error": lock_err}
    try:
        prs.save(pptx_path)
    except PermissionError as e:
        return {"status": "error",
                "error": f"Cannot save pptx — file is locked. Close PowerPoint and retry. ({e})"}

    return {
        "status": "ok",
        "month": month_date.isoformat(),
        "visits": {"actual": seo_visits, "target": visits_target, "pct": round(visits_pct, 1)},
        "revenue": {"actual": seo_omzet, "target": omzet_target, "pct": round(omzet_pct, 1)},
    }


def update_serp_table(pptx_path: str, excel_path: str, slide_index: int = 1,
                      target_yyyymm: Optional[int] = None) -> Dict:
    """Update the SERP rankings table on the given slide (0-based).

    When `target_yyyymm` is given, the table reports that month and the one
    before it; otherwise the two rightmost month columns are used."""
    lock_err = _check_pptx_lock(pptx_path)
    if lock_err:
        return {"status": "error", "error": lock_err}

    src = _read_serp_last_two_months(excel_path, target_yyyymm)
    if src is None:
        return {"status": "error", "error": "Could not read 2 month columns from `serp` sheet"}
    prev_name, last_name, prev_vals, last_vals = src

    prs = Presentation(pptx_path)
    if slide_index >= len(prs.slides):
        return {"status": "error",
                "error": f"Slide index {slide_index} out of range"}
    slide = prs.slides[slide_index]
    shape, tbl = _find_serp_table(slide)
    if tbl is None:
        return {"status": "error", "error": "SERP table not found on slide"}

    # Header (row 0): Type URL | prev_name | last_name | Delta
    _set_cell_text(tbl.cell(0, 1), prev_name)
    _set_cell_text(tbl.cell(0, 2), last_name)
    _set_cell_text(tbl.cell(0, 3), "Delta")

    # Data rows: assume the existing rows 1..4 are in URL_TYPE_ORDER, matched by label
    rows_written = []
    for r_idx in range(1, len(tbl.rows)):
        label = tbl.cell(r_idx, 0).text.strip()
        if label not in URL_TYPE_ORDER:
            continue
        prev_v = prev_vals.get(label)
        last_v = last_vals.get(label)
        _set_cell_text(tbl.cell(r_idx, 1), _fmt_position(prev_v))
        _set_cell_text(tbl.cell(r_idx, 2), _fmt_position(last_v))
        _set_cell_text(tbl.cell(r_idx, 3), _fmt_delta(prev_v, last_v))
        rows_written.append(label)

    lock_err = _check_pptx_lock(pptx_path)
    if lock_err:
        return {"status": "error", "error": lock_err}
    try:
        prs.save(pptx_path)
    except PermissionError as e:
        return {"status": "error",
                "error": f"Cannot save pptx — file is locked. Close PowerPoint and retry. ({e})"}

    return {
        "status": "ok",
        "table": shape.name,
        "months": [prev_name, last_name],
        "rows_written": rows_written,
    }
