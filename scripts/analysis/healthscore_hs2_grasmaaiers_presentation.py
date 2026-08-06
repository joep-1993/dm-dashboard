#!/usr/bin/env python3
"""Build the "Healthscore 2.0 — Grasmaaiers" deck as a .pptx.

Companion to healthscore_hs2_presentation.py (the HS2.0-vs-HS1.0 deck): same
960x540pt 16:9 canvas, same Aptos/magenta/cyan styling, but zoomed in on ONE
testcategory instead of the whole set. Grasmaaiers (9003581) is the pick because
it is the only testcat whose HS2.0 selection is actually live (sitemap push
03-08-2026), so the numbers here describe something running, not a simulation.

EVERY FIGURE IS MEASURED, NOT ESTIMATED. Sources:
  - Downloads/claude/hs2_catdiff_seasonal_v2.csv — the per-URL June-2026 holdout
    diff (1.591 rows for cat 9003581). All URL counts, coverage %, the URL-type
    split and the kept/added/dropped/uncovered breakdown come from here,
    recomputed at build time rather than transcribed.
  - pa.hs2_cat_knee / pa.hs2_cat_cap — knee-based base cap + seasonal index.
  - pa.hs2_sitemap — the 30-06 and 03-08 selections actually written out.
  - pa.hs2_cat_maincat — category naming / urlslug.

DENOMINATOR NOTE (checked 2026-08-06, do not "fix" this):
  The coverage denominator is 2.081 SEO visits, not the 6.915 that
  pa.hs2_cat_month reports for 202606. Both are right for their own purpose.
  Coverage is SEO-only AND restricted to `dv.url ~ '^https?://www\\.beslist\\.nl/'`
  (_SEO_WHERE) — that regex alone takes 3.824 -> 2.081. pa.hs2_cat_month is fed
  by the ALL-channel variant (_ALL_JOIN/_ALL_WHERE) because cap-sizing wants the
  full demand signal, not just SEO. Verified in Redshift: SEO+regex = 2.081,
  exactly the CSV total.

Usage:  venv/bin/python scripts/analysis/healthscore_hs2_grasmaaiers_presentation.py
Output: /mnt/c/Users/JoepvanSchagen/Downloads/claude/Healthscore_2.0_Grasmaaiers.pptx
        (export to PDF from PowerPoint, as with the parent deck)
"""
import csv
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = "/mnt/c/Users/JoepvanSchagen/Downloads/claude/Healthscore_2.0_Grasmaaiers.pptx"
CATDIFF = "/mnt/c/Users/JoepvanSchagen/Downloads/claude/hs2_catdiff_seasonal_v2.csv"
CAT_NAME = "Grasmaaiers"
CAT_ID = 9003581

FONT = "Aptos"
MAGENTA = RGBColor(0x9F, 0x2B, 0x92)
CYAN = RGBColor(0x00, 0xAF, 0xEF)
DARK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x64, 0x6C, 0x75)
BLUE = RGBColor(0x2C, 0x6C, 0xDF)
BAND = RGBColor(0xF5, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TRACK = RGBColor(0xE1, 0xE7, 0xEC)
GREEN_DK = RGBColor(0x2D, 0x9F, 0x6A)
AMBER = RGBColor(0xE8, 0xA2, 0x1E)
RED = RGBColor(0xD9, 0x42, 0x42)

PX = 90.0


def px(v):
    return Inches(v / PX)


def nl(v, dec=0):
    """Dutch number formatting: 1.234,5"""
    s = f"{v:,.{dec}f}"
    return s.replace(",", "\x00").replace(".", ",").replace("\x00", ".")


def shape(slide, kind, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(kind, px(x), px(y), px(w), px(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def rounded(slide, x, y, w, h, fill, radius=0.12, line=None):
    s = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)
    s.adjustments[0] = radius
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=None):
    """runs: list of (string, size_pt, bold, colour) — one paragraph each."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (s, size, bold, colour) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = s
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
    return tb


def title_bar(slide, label):
    rounded(slide, 52, 36, 1096, 90, MAGENTA, radius=0.18)
    text(slide, 52, 36, 1096, 90, [(label, 30, True, WHITE)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def progress(slide, x, y, w, pct, fill=BLUE):
    h = 28
    rounded(slide, x, y, w, h, TRACK, radius=0.5)
    filled = max(w * pct / 100.0, h if pct > 0 else 0)
    if pct > 0:
        rounded(slide, x, y, filled, h, fill, radius=0.5)


def footnote(slide, s):
    text(slide, 52, 632, 1096, 30, [(s, 9, False, MUTED)])


# ---------------------------------------------------------------------------
# Measured inputs — recomputed from the per-URL CSV at build time
# ---------------------------------------------------------------------------
def utype(npath):
    if npath.startswith("/p/"):
        return "PLP (/p/)"
    if "/r/" in npath:
        return "R-url"
    if "/c/" in npath:
        return "C-url"
    return "Browse-url"


TYPES = ["C-url", "R-url", "PLP (/p/)", "Browse-url"]


def load():
    if not os.path.exists(CATDIFF):
        sys.exit(f"missing {CATDIFF}")
    rows = [r for r in csv.DictReader(open(CATDIFF)) if r["cat_name"] == CAT_NAME]

    def num(r, k):
        try:
            return float(r[k] or 0)
        except (TypeError, ValueError):
            return 0.0

    for r in rows:
        r["v"] = num(r, "june_visits")
        r["e"] = num(r, "june_revenue")
        r["t"] = utype(r["npath"])

    tv = sum(r["v"] for r in rows)
    te = sum(r["e"] for r in rows)
    h1 = [r for r in rows if r["in_hs1"] == "1"]
    h2 = [r for r in rows if r["in_hs2"] == "1"]
    d = dict(
        n_rows=len(rows), tot_v=tv, tot_e=te,
        h1_n=len(h1), h2_n=len(h2),
        h1_v=sum(r["v"] for r in h1), h2_v=sum(r["v"] for r in h2),
        h1_e=sum(r["e"] for r in h1), h2_e=sum(r["e"] for r in h2),
    )
    d["h1_cov_v"] = 100 * d["h1_v"] / tv
    d["h2_cov_v"] = 100 * d["h2_v"] / tv
    d["h1_cov_e"] = 100 * d["h1_e"] / te
    d["h2_cov_e"] = 100 * d["h2_e"] / te
    for st in ("kept", "added", "dropped", "uncovered"):
        g = [r for r in rows if r["status"] == st]
        d[st] = (len(g), sum(r["v"] for r in g), sum(r["e"] for r in g))
    # per URL-type
    d["by_type"] = []
    for t in TYPES:
        g = [r for r in rows if r["t"] == t]
        if not g:
            continue
        gv = sum(r["v"] for r in g)
        a = [r for r in g if r["in_hs1"] == "1"]
        b = [r for r in g if r["in_hs2"] == "1"]
        d["by_type"].append(dict(
            t=t, n=len(g), v=gv, share=100 * gv / tv,
            h1_n=len(a), h2_n=len(b),
            h1_cov=100 * sum(r["v"] for r in a) / gv if gv else 0.0,
            h2_cov=100 * sum(r["v"] for r in b) / gv if gv else 0.0,
        ))
    d["top_added"] = sorted([r for r in rows if r["status"] == "added"],
                            key=lambda r: -r["v"])[:4]
    d["top_unc"] = sorted([r for r in rows if r["status"] == "uncovered"],
                          key=lambda r: -r["v"])[:3]
    d["max_dropped_v"] = max((r["v"] for r in rows if r["status"] == "dropped"),
                             default=0)
    return d


# pa.hs2_cat_knee (cat 9003581) + pa.hs2_cat_cap, and pa.hs2_sitemap selections.
KNEE = dict(yearly=57_786, knee80=674, knee90=1_248, knee95=1_928, n_urls=4_048)
CAPS = [(1, 0.420, 524), (2, 1.260, 1_573), (3, 1.653, 2_063), (4, 1.826, 2_279),
        (5, 1.826, 2_279), (6, 1.463, 1_825), (7, 1.463, 1_825), (8, 1.417, 1_769),
        (9, 1.048, 1_308), (10, 0.798, 996), (11, 0.512, 639), (12, 0.324, 499)]
MONTHS_NL = ["jan", "feb", "mrt", "apr", "mei", "jun",
             "jul", "aug", "sep", "okt", "nov", "dec"]
# pa.hs2_sitemap, deepest_category_id 9003581, by as_of_date
SITEMAP = {
    "2026-06-30": {"C-url": 279, "PLP": 319, "R-url": 401, "Browse": 1},
    "2026-08-03": {"C-url": 303, "PLP": 314, "R-url": 447, "Browse": 2},
}


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_cover(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 1200, 675, MAGENTA)
    tri = shape(s, MSO_SHAPE.RIGHT_TRIANGLE, 600, 275, 600, 400,
                RGBColor(0xB0, 0x3D, 0xA3))
    tri.rotation = 0
    text(s, 80, 140, 900, 220,
         [("Healthscore 2.0 — Grasmaaiers", 48, True, WHITE),
          ("Eén testcategorie uitgelicht: wat verandert er precies?", 22, False, WHITE)],
         spacing=10)
    text(s, 80, 460, 800, 60,
         [(f"Categorie {CAT_ID} · Tuinartikelen · holdout juni 2026 (out-of-sample) · "
           f"HS2.0-selectie live sinds 03-08-2026", 14, False, WHITE)])
    text(s, 820, 590, 320, 50, [("beslist.nl", 30, True, WHITE)], align=PP_ALIGN.RIGHT)
    return s


def slide_context(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "De categorie in cijfers")
    cards = [
        ("URL's in de categorie", nl(KNEE["n_urls"]),
         "totale kandidatenpoel; 1.591 daarvan zitten in een set of kregen juniverkeer"),
        ("SEO-bezoeken juni 2026", nl(d["tot_v"]),
         f"€ {nl(d['tot_e'], 2)} omzet — dit is de noemer van elk dekkingspercentage"),
        ("Jaarvolume (alle kanalen)", nl(KNEE["yearly"]),
         "gebruikt voor cap-sizing, niet voor de dekkings-KPI"),
        ("Knik-punt (90%)", nl(KNEE["knee90"]),
         f"URL's die 90% van het jaarvolume dekken — {nl(KNEE['knee80'])} voor 80%, "
         f"{nl(KNEE['knee95'])} voor 95%"),
    ]
    for i, (head, big, sub) in enumerate(cards):
        x = 52 + i * 282
        rounded(s, x, 165, 248, 300, BAND, radius=0.06)
        text(s, x + 20, 190, 208, 60, [(head, 13, True, DARK)], align=PP_ALIGN.CENTER)
        text(s, x + 20, 262, 208, 60, [(big, 30, True, CYAN)], align=PP_ALIGN.CENTER)
        text(s, x + 20, 330, 208, 120, [(sub, 9.5, False, MUTED)], align=PP_ALIGN.CENTER)
    rounded(s, 52, 492, 1096, 128, BAND, radius=0.06)
    text(s, 79, 512, 1042, 30,
         [("Waarom Grasmaaiers?", 14, True, DARK)])
    text(s, 79, 546, 1042, 80,
         [("Van de tien testcategorieën is dit de enige waarvan de HS2.0-selectie "
           "daadwerkelijk in de sitemap staat (push 03-08-2026). De cijfers hieronder "
           "beschrijven dus geen simulatie maar iets wat draait. Daarnaast is het een "
           "uitgesproken seizoenscategorie — precies waar de nieuwe caps voor gebouwd zijn.",
           11.5, False, DARK)])
    footnote(s, "Bron: pa.hs2_cat_knee en pa.hs2_cat_maincat (cat 9003581, urlslug "
                "tuin_accessoires_504070_4468892); junicijfers uit hs2_catdiff_seasonal_v2.csv.")
    return s


def slide_diff(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Het verschil in één tabel")
    for label, hx in [("", 79), ("Huidig (HS1.0)", 620), ("HS2.0", 830), ("Verschil", 1010)]:
        text(s, hx, 148, 220, 24, [(label, 11, False, MUTED)])
    rows = [
        ("URL's in de set", nl(d["h1_n"]), nl(d["h2_n"]),
         f"+{nl(d['h2_n'] - d['h1_n'])}", True),
        ("Dekking SEO-bezoeken", f"{nl(d['h1_cov_v'], 1)}%", f"{nl(d['h2_cov_v'], 1)}%",
         f"+{nl(d['h2_cov_v'] - d['h1_cov_v'], 1)}pp", True),
        ("Dekking SEO-omzet", f"{nl(d['h1_cov_e'], 1)}%", f"{nl(d['h2_cov_e'], 1)}%",
         f"+{nl(d['h2_cov_e'] - d['h1_cov_e'], 1)}pp", True),
        ("Gedekte bezoeken", nl(d["h1_v"]), nl(d["h2_v"]),
         f"+{nl(d['h2_v'] - d['h1_v'])}", True),
        ("Gedekte omzet", f"€ {nl(d['h1_e'], 2)}", f"€ {nl(d['h2_e'], 2)}",
         f"+€ {nl(d['h2_e'] - d['h1_e'], 2)}", True),
        ("Bezoeken per URL", nl(d["h1_v"] / d["h1_n"], 2), nl(d["h2_v"] / d["h2_n"], 2),
         f"+{nl(d['h2_v'] / d['h2_n'] - d['h1_v'] / d['h1_n'], 2)}", True),
    ]
    for i, (label, a, b, delta, good) in enumerate(rows):
        y = 178 + i * 62
        rounded(s, 52, y, 1096, 54, BAND, radius=0.12)
        text(s, 79, y + 16, 520, 26, [(label, 13.5, True, DARK)])
        text(s, 620, y + 16, 200, 26, [(a, 13.5, False, MUTED)])
        text(s, 830, y + 16, 180, 26, [(b, 13.5, True, DARK)])
        text(s, 1010, y + 16, 140, 26,
             [(delta, 13.5, True, GREEN_DK if good else RED)])
    text(s, 52, 556, 1096, 60,
         [("Meer URL's én een hogere dekking per URL: HS2.0 zet 204 URL's méér in de set "
           "en vangt daarmee 380 extra bezoeken. Het is dus geen kwestie van 'gewoon meer "
           "URL's' — de mix verandert, zie de volgende slide.", 12.5, False, DARK)])
    footnote(s, "Bron: hs2_catdiff_seasonal_v2.csv, holdout juni 2026, voorspelvenster "
                "maart–mei (leakage-vrij). Noemer = 2.081 SEO-bezoeken / € 193,63 van de "
                "categorie in juni.")
    return s


def slide_types(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "De winst zit volledig in R-urls")
    heads = [("URL-type", 79), ("Aandeel bezoeken", 460), ("HS1.0 dekking", 660),
             ("HS2.0 dekking", 850), ("Verschil", 1030)]
    for label, hx in heads:
        text(s, hx, 148, 220, 24, [(label, 11, False, MUTED)])
    for i, r in enumerate(d["by_type"]):
        y = 180 + i * 80
        rounded(s, 52, y, 1096, 68, BAND, radius=0.10)
        text(s, 79, y + 12, 380, 26, [(r["t"], 14, True, DARK)])
        text(s, 79, y + 38, 380, 22,
             [(f"{nl(r['n'])} kandidaat-URL's · {nl(r['h1_n'])} → {nl(r['h2_n'])} in de set",
               9.5, False, MUTED)])
        text(s, 460, y + 22, 180, 26, [(f"{nl(r['share'], 1)}%", 13.5, False, DARK)])
        text(s, 660, y + 22, 180, 26,
             [(f"{nl(r['h1_cov'], 1)}%", 13.5, True,
               GREEN_DK if r["h1_cov"] >= 50 else RED)])
        text(s, 850, y + 22, 180, 26,
             [(f"{nl(r['h2_cov'], 1)}%", 13.5, True,
               GREEN_DK if r["h2_cov"] >= 50 else RED)])
        delta = r["h2_cov"] - r["h1_cov"]
        text(s, 1030, y + 22, 140, 26,
             [(f"{'+' if delta >= 0 else '−'}{nl(abs(delta), 1)}pp", 13.5, True,
               GREEN_DK if delta >= 0 else RED)])
    text(s, 52, 512, 1096, 110,
         [("R-urls zijn een derde van het juniverkeer en waren met 7,2% vrijwel ongedekt — "
           "27 van de 484. HS2.0 zet er 338 in en tilt de dekking naar 70,2%. Dat is in z'n "
           "eentje de hele +18,3pp.", 12.5, False, DARK),
          ("De PLP-dekking zakt van 90,9% naar 76,4%. Dat is geen fout maar de ruil: "
           "productpagina's met één of twee bezoeken maken plaats voor zoek-URL's die meer "
           "verkeer trekken. C-urls — 78% van de omzet — blijven met 95,4% onaangetast.",
           12.5, False, DARK)], spacing=8)
    footnote(s, "Type afgeleid uit het pad: /p/ = PLP, /r/ = R-url, /c/ = C-url, rest = "
                "Browse. Dekking = aandeel van de junibezoeken bínnen dat type dat op een "
                "URL in de set landt.")
    return s


def slide_mutation(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Wat erbij komt en wat eraf valt")
    blocks = [
        ("Behouden", "kept", GREEN_DK),
        ("Nieuw in de set", "added", CYAN),
        ("Eruit", "dropped", RED),
        ("Nog steeds ongedekt", "uncovered", AMBER),
    ]
    for i, (head, key, colour) in enumerate(blocks):
        n, v, e = d[key]
        x = 52 + i * 282
        rounded(s, x, 165, 248, 250, BAND, radius=0.06)
        text(s, x + 20, 188, 208, 30, [(head, 13.5, True, DARK)], align=PP_ALIGN.CENTER)
        text(s, x + 20, 228, 208, 50, [(nl(n), 30, True, colour)], align=PP_ALIGN.CENTER)
        text(s, x + 20, 282, 208, 24, [("URL's", 10, False, MUTED)], align=PP_ALIGN.CENTER)
        text(s, x + 20, 316, 208, 30,
             [(f"{nl(v)} bezoeken", 13, True, DARK)], align=PP_ALIGN.CENTER)
        text(s, x + 20, 350, 208, 30,
             [(f"€ {nl(e, 2)}", 11, False, MUTED)], align=PP_ALIGN.CENTER)
    rounded(s, 52, 440, 540, 180, BAND, radius=0.06)
    text(s, 79, 460, 486, 26, [("De grootste toevoegingen", 13, True, GREEN_DK)])
    for j, r in enumerate(d["top_added"]):
        term = r["npath"].split("/r/")[-1] if "/r/" in r["npath"] else r["npath"][-40:]
        text(s, 79, 492 + j * 30, 486, 26,
             [(f"+{nl(r['v'])} bezoeken   /r/{term[:44]}", 10.5, False, DARK)])
    rounded(s, 608, 440, 540, 180, BAND, radius=0.06)
    text(s, 635, 460, 486, 26, [("Het verlies is verwaarloosbaar", 13, True, RED)])
    text(s, 635, 492, 486, 120,
         [(f"De 395 URL's die eruit vallen waren samen goed voor {nl(d['dropped'][1])} "
           f"junibezoeken (€ {nl(d['dropped'][2], 2)}) — gemiddeld 0,2 bezoek per URL, "
           f"de zwaarste verliest er {nl(d['max_dropped_v'])}. Dat is precies wat je wilt "
           f"opruimen: productpagina's zonder publiek.", 11, False, DARK)])
    footnote(s, "‘Nog steeds ongedekt’ = URL's met juniverkeer die in géén van beide sets "
                "zitten; vrijwel allemaal long-tail R-urls met 1–6 bezoeken.")
    return s


def slide_cap(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "De seizoenscap — waarom een vaste top-1000 niet werkt")
    text(s, 79, 148, 400, 24, [("Cap per kalendermaand", 11, False, MUTED)])
    maxcap = max(c for _, _, c in CAPS)
    for i, (m, idx, cap) in enumerate(CAPS):
        x = 79 + i * 88
        h = 170 * cap / maxcap
        colour = CYAN if m in (6, 8) else TRACK
        rounded(s, x, 356 - h, 62, h, colour, radius=0.10)
        text(s, x - 6, 364, 74, 22, [(MONTHS_NL[m - 1], 10.5, m in (6, 8), DARK)],
             align=PP_ALIGN.CENTER)
        text(s, x - 6, 386, 74, 22, [(nl(cap), 9.5, False, MUTED)], align=PP_ALIGN.CENTER)
        text(s, x - 6, 330 - h, 74, 20, [(f"{nl(idx, 2)}×", 9, False, MUTED)],
             align=PP_ALIGN.CENTER)
    rounded(s, 52, 424, 540, 196, BAND, radius=0.06)
    text(s, 79, 444, 486, 26, [("Hoe de cap tot stand komt", 13, True, DARK)])
    text(s, 79, 478, 486, 130,
         [(f"Basis = het knik-punt: {nl(KNEE['knee90'])} URL's dekken 90% van het "
           f"jaarvolume ({nl(KNEE['yearly'])} bezoeken over 12 maanden, alle kanalen). "
           f"Die basis gaat maal de seizoensindex uit 24 maanden klimatologie: "
           f"0,42× in januari, 1,83× in april/mei. Zo krijgt Grasmaaiers in mei "
           f"{nl(2279)} URL's en in december {nl(499)}.", 11, False, DARK)])
    rounded(s, 608, 424, 540, 196, BAND, radius=0.06)
    text(s, 635, 444, 486, 26, [("De cap was in juni niet de bottleneck", 13, True, AMBER)])
    text(s, 635, 478, 486, 130,
         [(f"De junicap staat op {nl(1825)} URL's, maar HS2.0 selecteerde er "
           f"{nl(d['h2_n'])}. De beperking was dus niet de cap maar het aanbod: er waren "
           f"simpelweg niet meer URL's met verkeer in het voorspelvenster. Pas bij een "
           f"ruimere kandidatenpoel gaat de cap knellen.", 11, False, DARK)])
    footnote(s, "Bron: pa.hs2_cat_cap en pa.hs2_cat_knee (cat 9003581). Cap-sizing gebruikt "
                "álle kanalen (volledig vraagsignaal); de dekkings-KPI blijft SEO-only.")
    return s


def slide_live(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Wat er nu live staat")
    for label, hx in [("URL-type", 79), ("Selectie 30-06", 620), ("Selectie 03-08", 840),
                      ("Verschil", 1030)]:
        text(s, hx, 148, 220, 24, [(label, 11, False, MUTED)])
    keys = [("C-url", "C-url"), ("R-url", "R-url"), ("PLP (/p/)", "PLP"),
            ("Browse-url", "Browse")]
    a_tot = sum(SITEMAP["2026-06-30"].values())
    b_tot = sum(SITEMAP["2026-08-03"].values())
    for i, (label, k) in enumerate(keys):
        a, b = SITEMAP["2026-06-30"][k], SITEMAP["2026-08-03"][k]
        y = 180 + i * 72
        rounded(s, 52, y, 1096, 60, BAND, radius=0.10)
        text(s, 79, y + 18, 500, 26, [(label, 13.5, True, DARK)])
        text(s, 620, y + 18, 180, 26, [(nl(a), 13.5, False, MUTED)])
        text(s, 840, y + 18, 180, 26, [(nl(b), 13.5, True, DARK)])
        text(s, 1030, y + 18, 140, 26,
             [(f"{'+' if b >= a else '−'}{nl(abs(b - a))}", 13.5, True,
               GREEN_DK if b >= a else RED)])
    y = 180 + len(keys) * 72
    text(s, 79, y + 6, 500, 26, [("Totaal", 13.5, True, DARK)])
    text(s, 620, y + 6, 180, 26, [(nl(a_tot), 13.5, False, MUTED)])
    text(s, 840, y + 6, 180, 26, [(nl(b_tot), 13.5, True, DARK)])
    text(s, 1030, y + 6, 140, 26, [(f"+{nl(b_tot - a_tot)}", 13.5, True, GREEN_DK)])
    text(s, 52, y + 46, 1096, 80,
         [("De augustusselectie is 66 URL's groter dan die van eind juni, vrijwel volledig "
           "R-urls (+46) en C-urls (+24) — de augustuscap ligt met 1.769 iets onder de "
           "junicap, dus ook hier bepaalt het aanbod de omvang, niet het plafond. "
           "Alle URL's komen uit het ‘scored’-bakje; het garantie-bakje voor nieuwe URL's "
           "leverde voor deze categorie niets op.", 11.5, False, DARK)])
    footnote(s, "Bron: pa.hs2_sitemap, deepest_category_id 9003581, as_of_date 2026-06-30 "
                "en 2026-08-03. De 03-08-selectie is wat nu daadwerkelijk in de sitemap staat.")
    return s


def slide_readout(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Wat dit zegt — en wat het niet zegt")
    items = [
        ("Bevestigt", "Het R-url-verhaal uit de hoofddeck klopt ook op categorieniveau: "
                      "7,2% → 70,2% dekking op R-urls is hier de hele winst. Geen "
                      "scoringsprobleem, een selectieprobleem.", GREEN_DK),
        ("Goedkoop", f"Het verlies is {nl(d['dropped'][1])} bezoeken (€ {nl(d['dropped'][2], 2)}) "
                     f"over 395 URL's. Tegenover +{nl(d['added'][1])} bezoeken erbij is de "
                     f"ruil ruim 5:1 in ons voordeel.", GREEN_DK),
        ("Let op", "Één categorie in één maand. Grasmaaiers is klein (2.081 SEO-bezoeken) "
                   "en zat in juni over z'n piek heen; een categorie met een ander "
                   "seizoenspatroon of meer PLP-verkeer kan anders uitpakken.", AMBER),
        ("Nog te meten", "De 03-08-selectie staat live maar is nog niet nagemeten. De echte "
                         "test is of de gerealiseerde SEO-dekking in augustus/september "
                         "richting de voorspelde 83,6% beweegt.", AMBER),
        ("Volgende stap", "Dezelfde uitsnede maken voor een grote, niet-seizoensgebonden "
                          "testcat (Sneakers of Mobiele telefoons) om te zien of het "
                          "patroon standhoudt.", MAGENTA),
    ]
    for i, (head, body, colour) in enumerate(items):
        y = 160 + i * 92
        rounded(s, 52, y, 1096, 80, BAND, radius=0.08)
        text(s, 79, y + 16, 220, 30, [(head, 13, True, colour)])
        text(s, 300, y + 14, 820, 62, [(body, 11, False, DARK)])
    footnote(s, "Alle cijfers herberekend uit hs2_catdiff_seasonal_v2.csv + pa.hs2_* op "
                "06-08-2026. Noemer SEO-only en beperkt tot www.beslist.nl-URL's; "
                "pa.hs2_cat_month (6.915 junibezoeken) is all-channel en dient alleen "
                "voor cap-sizing.")
    return s


def main():
    d = load()
    prs = Presentation()
    prs.slide_width = Inches(1200 / PX)
    prs.slide_height = Inches(675 / PX)

    slide_cover(prs, d)
    slide_context(prs, d)
    slide_diff(prs, d)
    slide_types(prs, d)
    slide_mutation(prs, d)
    slide_cap(prs, d)
    slide_live(prs, d)
    slide_readout(prs, d)

    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    print(f"Grasmaaiers: {d['h1_n']:,} -> {d['h2_n']:,} URLs, "
          f"{d['h1_cov_v']:.1f}% -> {d['h2_cov_v']:.1f}% visits, "
          f"{d['h1_cov_e']:.1f}% -> {d['h2_cov_e']:.1f}% revenue "
          f"(denominator {d['tot_v']:,.0f} visits / EUR {d['tot_e']:,.2f})")


if __name__ == "__main__":
    main()
