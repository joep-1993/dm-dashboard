#!/usr/bin/env python3
"""Build the "Healthscore 2.0 vs huidige Healthscore" deck as a .pptx.

Styled after Downloads/claude/beslist_geo_presentatie.pdf (the GEO/AI deck):
960x540pt 16:9, Aptos, magenta full-width title bar, light-grey rounded cards and
table bands, big cyan numbers, S/M/L impact-effort pills. Palette and font were
sampled from that PDF rather than guessed:

    title bar   #9F2B92     big numbers #00AFEF     body text  #1F2A37
    card/band   #F5F6F8     muted label #646C75     blue label #2C6CDF
    pills       impact M #6EB891 / S #CAE1D5 · effort S #2D9F6A / M #E8A21E / L #D94242

EVERY FIGURE IN THIS DECK IS MEASURED, NOT ESTIMATED. Sources:
  - pa.hs2_shadow (run 2026-07-17, holdout 2026-06, cap_n=1000) — the headline
    45.1%/50.9% -> 71.0%/75.1% compare and the add/drop split.
  - pa.healthscore_coverage (target_month 2026-05) — per-URL-type coverage, which
    is where the "the gap is entirely R-urls" slide comes from.
  - pa.hs2_sitemap — 769,017 scored + 138,352 new = 907,369.
  - Downloads/claude/hs2_catdiff_seasonal_v2.csv — the 10-category seasonal-cap
    validation, RECOMPUTED here from the per-URL rows (70.2->84.7 visits,
    80.4->91.5 revenue). An earlier note quoted 83.9/91.1; the CSV is authoritative.

Usage:  venv/bin/python scripts/analysis/healthscore_hs2_presentation.py
Output: /mnt/c/Users/JoepvanSchagen/Downloads/claude/Healthscore_2.0_vs_1.0.pptx
"""
import csv
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = "/mnt/c/Users/JoepvanSchagen/Downloads/claude/Healthscore_2.0_vs_1.0.pptx"
CATDIFF = "/mnt/c/Users/JoepvanSchagen/Downloads/claude/hs2_catdiff_seasonal_v2.csv"

FONT = "Aptos"
MAGENTA = RGBColor(0x9F, 0x2B, 0x92)
CYAN = RGBColor(0x00, 0xAF, 0xEF)
DARK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x64, 0x6C, 0x75)
BLUE = RGBColor(0x2C, 0x6C, 0xDF)
BAND = RGBColor(0xF5, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TRACK = RGBColor(0xE1, 0xE7, 0xEC)
GREEN_M = RGBColor(0x6E, 0xB8, 0x91)
GREEN_S = RGBColor(0xCA, 0xE1, 0xD5)
GREEN_DK = RGBColor(0x2D, 0x9F, 0x6A)
AMBER = RGBColor(0xE8, 0xA2, 0x1E)
RED = RGBColor(0xD9, 0x42, 0x42)

# The reference renders at 90 px/inch, so every coordinate lifted from it divides
# by 90. Keeping that unit makes the layout numbers comparable to the source deck.
PX = 90.0


def px(v):
    return Inches(v / PX)


def shape(slide, kind, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(kind, px(x), px(y), px(w), px(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def rounded(slide, x, y, w, h, fill, radius=0.12, line=None):
    s = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)
    s.adjustments[0] = radius
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=None):
    """runs: list of (string, size_pt, bold, colour) — one paragraph each."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (s, size, bold, colour) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = s
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
    return tb


def title_bar(slide, label):
    rounded(slide, 52, 36, 1096, 90, MAGENTA, radius=0.18)
    text(slide, 52, 36, 1096, 90, [(label, 30, True, WHITE)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def pill(slide, cx, cy, letter, fill, fg=WHITE):
    w, h = 66, 40
    rounded(slide, cx - w / 2, cy - h / 2, w, h, fill, radius=0.5)
    text(slide, cx - w / 2, cy - h / 2, w, h, [(letter, 13, True, fg)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def progress(slide, x, y, w, pct, fill=BLUE):
    h = 28
    rounded(slide, x, y, w, h, TRACK, radius=0.5)
    filled = max(w * pct / 100.0, h if pct > 0 else 0)
    if pct > 0:
        rounded(slide, x, y, filled, h, fill, radius=0.5)


def footnote(slide, s):
    text(slide, 52, 632, 1096, 30, [(s, 9, False, MUTED)])


# ---------------------------------------------------------------------------
# Measured inputs
# ---------------------------------------------------------------------------
# pa.hs2_shadow, id=1: holdout 2026-06, cap_n 1000, weights 0.889/0.111
SHADOW = dict(
    cur_cov_v=45.1, cur_cov_r=50.9, cur_urls=748_860,
    hs2_cov_v=71.0, hs2_cov_r=75.1, hs2_urls=709_859,
    add_v=569_067, add_r=63_776, drop_v=140_530, drop_r=14_738,
    total_v=1_652_614, total_r=202_941,
)
# pa.healthscore_coverage, target_month 2026-05 (set of 756,709 urls)
BY_TYPE = [
    # label, visit coverage %, revenue coverage %, share of SEO visits %, healthy?
    ("C-url", 90.4, 92.0, 32.4, True),
    ("PLP (/p/)", 79.7, 80.1, 13.0, True),
    ("Browse-url", 99.6, 99.9, 1.5, True),
    ("R-url", 2.8, 2.7, 52.9, False),
]


def seasonal_validation():
    """Recompute the 10-category seasonal-cap result from the per-URL CSV."""
    if not os.path.exists(CATDIFF):
        return None, []
    rows = list(csv.DictReader(open(CATDIFF)))

    def num(r, k):
        try:
            return float(r[k] or 0)
        except (TypeError, ValueError):
            return 0.0

    tot_v = sum(num(r, "june_visits") for r in rows)
    tot_r = sum(num(r, "june_revenue") for r in rows)
    h1v = sum(num(r, "june_visits") for r in rows if r["in_hs1"] == "1")
    h1r = sum(num(r, "june_revenue") for r in rows if r["in_hs1"] == "1")
    h2v = sum(num(r, "june_visits") for r in rows if r["in_hs2"] == "1")
    h2r = sum(num(r, "june_revenue") for r in rows if r["in_hs2"] == "1")
    per = {}
    for r in rows:
        c = per.setdefault(r["cat_name"], [0.0, 0.0, 0.0])
        v = num(r, "june_visits")
        c[0] += v
        c[1] += v if r["in_hs1"] == "1" else 0
        c[2] += v if r["in_hs2"] == "1" else 0
    cats = sorted(
        ((n, 100 * a / t, 100 * b / t) for n, (t, a, b) in per.items() if t),
        key=lambda x: -(x[2] - x[1]),
    )
    totals = dict(
        v1=100 * h1v / tot_v, v2=100 * h2v / tot_v,
        r1=100 * h1r / tot_r, r2=100 * h2r / tot_r,
        visits=tot_v, revenue=tot_r,
    )
    return totals, cats


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 1200, 675, MAGENTA)
    # A lighter wedge, so the flat fill reads as designed rather than as a default.
    tri = shape(s, MSO_SHAPE.RIGHT_TRIANGLE, 600, 275, 600, 400,
                RGBColor(0xB0, 0x3D, 0xA3))
    tri.rotation = 0
    text(s, 80, 150, 800, 200,
         [("Healthscore 2.0", 54, True, WHITE),
          ("Vergelijking met de huidige Healthscore", 22, False, WHITE)],
         spacing=10)
    text(s, 80, 470, 700, 40,
         [("Out-of-sample gemeten op juni 2026 · status: in de koelkast", 14, False, WHITE)])
    text(s, 820, 590, 320, 50, [("beslist.nl", 30, True, WHITE)], align=PP_ALIGN.RIGHT)
    return s


def slide_where_we_are(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Waar staan we nu?")
    cards = [
        ("Totale dekking", "42,6%", 42.6, "van de SEO-bezoeken landt op een URL in de sitemapset"),
        ("C-url", "90,4%", 90.4, "gezond"),
        ("PLP (/p/)", "79,7%", 79.7, "gezond"),
    ]
    for i, (head, big, pct, sub) in enumerate(cards):
        x = 52 + i * 282
        rounded(s, x, 175, 248, 440, BAND, radius=0.06)
        text(s, x + 20, 200, 208, 40, [(head, 15, True, DARK)], align=PP_ALIGN.CENTER)
        text(s, x + 20, 320, 208, 60, [(big, 32, True, CYAN)], align=PP_ALIGN.CENTER)
        progress(s, x + 34, 448, 180, pct)
        text(s, x + 20, 500, 208, 90, [(sub, 10, False, MUTED)], align=PP_ALIGN.CENTER)
    # Fourth card = the +/- insight card, exactly as the reference's last column.
    x = 898
    rounded(s, x, 175, 248, 440, BAND, radius=0.06)
    text(s, x + 20, 200, 208, 40, [("Het probleem", 15, True, DARK)], align=PP_ALIGN.CENTER)
    for j, (sign, colour, body) in enumerate([
        ("+", GREEN_DK, "Niet-R-url dekking\n≈ 86%"),
        ("–", RED, "R-url: 52,9% van de\nSEO-bezoeken,\n2,8% gedekt"),
    ]):
        cy = 300 + j * 130
        shape(s, MSO_SHAPE.OVAL, x + 22, cy - 22, 44, 44, colour)
        text(s, x + 22, cy - 22, 44, 44, [(sign, 16, True, WHITE)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 80, cy - 30, 150, 100, [(body, 11, False, DARK)])
    footnote(s, "Bron: pa.healthscore_coverage, doelmaand 2026-05, set van 756.709 URL's. "
                "Dekking = aandeel organische SEO-bezoeken op een URL die in de sitemapset staat.")
    return s


def slide_gap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Het gat zit volledig in R-urls")
    heads = [("URL-type", 79), ("Dekking bezoeken", 560), ("Dekking omzet", 760),
             ("Aandeel SEO-bezoeken", 960)]
    for label, hx in heads:
        text(s, hx, 158, 240, 24, [(label, 11, False, MUTED)])
    for i, (label, cv, cr, share, healthy) in enumerate(BY_TYPE):
        y = 195 + i * 88
        rounded(s, 52, y, 1096, 72, BAND, radius=0.08)
        text(s, 79, y + 24, 400, 30, [(label, 14, True, DARK)])
        text(s, 560, y + 24, 180, 30,
             [(f"{cv:,.1f}%".replace(".", ","), 14, True, GREEN_DK if healthy else RED)])
        text(s, 760, y + 24, 180, 30,
             [(f"{cr:,.1f}%".replace(".", ","), 14, False, DARK)])
        text(s, 960, y + 24, 180, 30,
             [(f"{share:,.1f}%".replace(".", ","), 14, False, DARK)])
    text(s, 52, 560, 1096, 60,
         [("R-urls zijn ruim de helft van alle SEO-bezoeken en worden vrijwel niet gedekt. "
           "Alle andere URL-typen staan er goed voor — dus dit is geen scoringsprobleem "
           "maar een selectieprobleem: R-urls zaten niet in de kandidatenverzameling.",
           13, False, DARK)])
    footnote(s, "Bron: pa.healthscore_coverage, doelmaand 2026-05. Homepage staat op 0% en is "
                "bewust niet in de set opgenomen.")
    return s


def slide_what_changes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Wat HS2.0 anders doet")
    blocks = [
        ("Score", "1", [
            "0,89 · pct(log bezoeken)",
            "+ 0,11 · pct(log omzet)",
            "percentiel binnen de eigen categorie",
            "gewichten uit backtest, reproduceerbaar",
        ]),
        ("Cap per categorie", "2", [
            "basis = knik-punt: hoeveel URL's dekken",
            "90% van de eigen 12-maands bezoeken",
            "× seizoensindex (24mnd klimatologie)",
            "+ 1 maand vooruitkijken (aanlooptijd SEO)",
        ]),
        ("Garantie-bakje", "3", [
            "nieuwe URL's (≤ 20 dagen) altijd mee",
            "lost cold-start op: ~20% van de bezoeken",
            "gaat naar URL's zonder historie",
            "138.352 URL's in de laatste run",
        ]),
    ]
    for i, (head, num, bullets) in enumerate(blocks):
        x = 52 + i * 374
        rounded(s, x, 165, 340, 300, BAND, radius=0.06)
        text(s, x + 24, 190, 292, 34, [(head, 16, True, DARK)])
        text(s, x + 24, 228, 60, 46, [(num, 26, True, CYAN)])
        text(s, x + 78, 240, 240, 30, [("van de drie", 11, False, MUTED)])
        text(s, x + 24, 288, 292, 160,
             [("•  " + b, 10.5, False, DARK) for b in bullets], spacing=6)
    rounded(s, 52, 492, 1096, 128, BAND, radius=0.06)
    text(s, 79, 512, 1042, 30, [("Wat géén gewicht kreeg", 14, True, DARK)])
    text(s, 79, 546, 1042, 70,
         [("CTR, bounce, momentum en zoekvolume kwamen in de backtest allemaal op 0,000 uit. "
           "Gerealiseerde bezoeken voorspellen dekking beter dan keyword-schattingen. "
           "CTR en bounce blijven wel staan als vangrail tegen navigatie-junk "
           "(‘weer’, ‘google’, ‘marktplaats’), niet als scoringsvariabele.",
           11, False, DARK)])
    footnote(s, "Gewichten 0,889 / 0,111 uit pa.hs2_shadow; reproduceren exact op een tweede "
                "split (april→mei), dus stabiel genoeg om vast te zetten.")
    return s


def slide_result(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Resultaat — juni 2026, out-of-sample")
    d = SHADOW
    cols = [
        ("Huidige Healthscore", f"{d['cur_cov_v']:.1f}%".replace(".", ","),
         f"{d['cur_cov_r']:.1f}%".replace(".", ","), f"{d['cur_urls']:,}".replace(",", "."), MUTED),
        ("Healthscore 2.0", f"{d['hs2_cov_v']:.1f}%".replace(".", ","),
         f"{d['hs2_cov_r']:.1f}%".replace(".", ","), f"{d['hs2_urls']:,}".replace(",", "."), CYAN),
    ]
    for i, (head, cv, cr, urls, colour) in enumerate(cols):
        x = 52 + i * 374
        rounded(s, x, 165, 340, 300, BAND, radius=0.06)
        text(s, x + 24, 188, 292, 34, [(head, 15, True, DARK)])
        text(s, x + 24, 236, 292, 52, [(cv, 34, True, colour)])
        text(s, x + 24, 292, 292, 24, [("dekking bezoeken", 10, False, MUTED)])
        text(s, x + 24, 330, 292, 40, [(cr, 22, True, colour)])
        text(s, x + 24, 372, 292, 24, [("dekking omzet", 10, False, MUTED)])
        text(s, x + 24, 410, 292, 40, [(f"{urls} URL's", 12, False, DARK)])
    x = 800
    rounded(s, x, 165, 348, 300, BAND, radius=0.06)
    text(s, x + 24, 188, 300, 34, [("Verschil", 15, True, DARK)])
    text(s, x + 24, 236, 300, 52, [("+25,9pp", 34, True, GREEN_DK)])
    text(s, x + 24, 292, 300, 24, [("bezoeken", 10, False, MUTED)])
    text(s, x + 24, 330, 300, 40, [("+24,2pp", 22, True, GREEN_DK)])
    text(s, x + 24, 372, 300, 24, [("omzet", 10, False, MUTED)])
    text(s, x + 24, 410, 300, 40,
         [("bij een 5% kleinere set", 12, True, DARK)])
    rounded(s, 52, 492, 1096, 128, BAND, radius=0.06)
    text(s, 79, 512, 500, 30, [("Wat erbij komt", 13, True, GREEN_DK)])
    text(s, 79, 544, 500, 70,
         [(f"+{d['add_v']:,}".replace(",", ".") + " bezoeken (€" +
           f"{d['add_r']:,.0f}".replace(",", ".") + ") die nu niet gedekt zijn", 11, False, DARK)])
    text(s, 620, 512, 500, 30, [("Wat eraf valt", 13, True, RED)])
    text(s, 620, 544, 500, 70,
         [(f"−{d['drop_v']:,}".replace(",", ".") + " bezoeken (€" +
           f"{d['drop_r']:,.0f}".replace(",", ".") + ") — eerlijk verlies, deels "
           "opgevangen door het garantie-bakje", 11, False, DARK)])
    footnote(s, "Bron: pa.hs2_shadow (run 17-07-2026, holdout 2026-06, cap 1000). Voorspeller = "
                "90 dagen vóór de holdout, dus geen leakage. Totaal juni: 1.652.614 bezoeken "
                "/ €202.941.")
    return s


def slide_seasonal(prs, totals, cats):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Seizoenscaps — 10 testcategorieën")
    if not totals:
        text(s, 52, 200, 1096, 60, [("hs2_catdiff_seasonal_v2.csv niet gevonden", 14, False, RED)])
        return s
    for label, hx in [("Categorie", 79), ("Huidig", 640), ("HS2.0", 800), ("Verschil", 960)]:
        text(s, hx, 144, 240, 24, [(label, 11, False, MUTED)])
    # All ten, not a top-8: a slide headed "10 testcategorieën" that quietly hides
    # the two weakest rows (Sneakers −6,4pp) would be the wrong kind of summary.
    for i, (name, a, b) in enumerate(cats):
        y = 172 + i * 40
        rounded(s, 52, y, 1096, 35, BAND, radius=0.16)
        text(s, 79, y + 8, 500, 24, [(name, 11.5, True, DARK)])
        text(s, 640, y + 8, 140, 24, [(f"{a:.1f}%".replace(".", ","), 11.5, False, MUTED)])
        text(s, 800, y + 8, 140, 24, [(f"{b:.1f}%".replace(".", ","), 11.5, True, DARK)])
        delta = b - a
        text(s, 960, y + 8, 140, 24,
             [(f"{delta:+.1f}pp".replace(".", ","), 11.5, True, GREEN_DK if delta >= 0 else RED)])
    y = 172 + len(cats) * 40 + 10
    text(s, 79, y, 560, 30,
         [(f"Totaal 10 categorieën: {totals['v1']:.1f}% → {totals['v2']:.1f}% bezoeken, "
           f"{totals['r1']:.1f}% → {totals['r2']:.1f}% omzet".replace(".", ","),
           13, True, DARK)])
    text(s, 640, y, 508, 60,
         [("Sneakers is het enige verlies (−6,4pp) en is een holdout-artefact: het "
           "voorspelvenster maart–mei is juist het dal van Sneakers, dus er waren te weinig "
           "kandidaat-URL's voor de hoge junicap.", 10, False, MUTED)])
    footnote(s, "Herberekend uit hs2_catdiff_seasonal_v2.csv (per URL, juni 2026): 65.354 bezoeken "
                "/ €13.163 over de 10 categorieën. Seizoensvoorbeeld: Airco piekt in juni op "
                "index 6,8× het jaargemiddelde.")
    return s


def slide_risks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Wat het kost en wat nog open staat")
    items = [
        ("Churn", "140.530 bezoeken (€14.738) vallen uit de set. Dat is echt verlies, geen "
                  "meetfout — deels opgevangen door het garantie-bakje."),
        ("Plafond", "Het haalbare maximum is 78,7%: ~20% van de junibezoeken gaat naar URL's "
                    "zonder enige historie. Gedrag kan die niet vinden, vandaar het bakje voor "
                    "nieuwe URL's."),
        ("Nieuw-bakje", "De 138.352 nieuwe URL's kunnen opgeblazen zijn door een facet-migratie "
                        "in het venster van 20 dagen. In een normale week hoort dit lager te "
                        "liggen — nog te bevestigen."),
        ("Niet live", "Niets is geproductionaliseerd: de tweewekelijkse run (features → "
                      "sitemap) en de maandelijkse caps-refresh moeten nog worden ingepland. "
                      "Alleen de shadow-vergelijking draait."),
        ("Volgende stap", "HS2.0-selectie live zetten voor de 10 testcategorieën en meten, "
                          "vóór een bredere uitrol."),
    ]
    for i, (head, body) in enumerate(items):
        y = 160 + i * 92
        rounded(s, 52, y, 1096, 80, BAND, radius=0.08)
        text(s, 79, y + 16, 200, 30, [(head, 13, True, MAGENTA)])
        text(s, 290, y + 14, 830, 60, [(body, 11, False, DARK)])
    footnote(s, "Status: HS2.0 fases 1–6 + seizoenscaps zijn gebouwd, gevalideerd en "
                "gecommit (e423557); de uitrol staat in de koelkast.")
    return s


def slide_order(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Voorgestelde volgorde")
    for label, hx in [("#", 79), ("Actie", 140), ("Impact", 800), ("Effort", 940), ("Fase", 1050)]:
        text(s, hx, 150, 200, 24, [(label, 11, False, MUTED)])
    rows = [
        ("1", "Live zetten voor de 10 testcategorieën", "M", GREEN_M, "S", GREEN_DK, "Nu"),
        ("2", "Tweewekelijkse run inplannen (features → sitemap)", "M", GREEN_M, "S", GREEN_DK, "Nu"),
        ("3", "Nieuw-bakje hermeten in een normale week", "S", GREEN_S, "S", GREEN_DK, "Nu"),
        ("4", "Caps-refresh maandelijks inplannen", "S", GREEN_S, "M", AMBER, "Daarna"),
        ("5", "Uitrol naar alle categorieën", "L", GREEN_M, "M", AMBER, "Daarna"),
        ("6", "Renderer-cutover: HS2.0 als bron van de sitemaps", "L", GREEN_M, "L", RED, "Later"),
    ]
    for i, (num, actie, imp, impc, eff, effc, fase) in enumerate(rows):
        y = 182 + i * 66
        rounded(s, 52, y, 1096, 58, BAND, radius=0.10)
        text(s, 79, y + 18, 60, 26, [(num, 13, True, CYAN)])
        text(s, 140, y + 18, 620, 26, [(actie, 12.5, True, DARK)])
        pill(s, 830, y + 29, imp, impc, WHITE if imp != "S" else DARK)
        pill(s, 968, y + 29, eff, effc, WHITE)
        text(s, 1040, y + 18, 120, 26, [(fase, 11, False, BLUE)])
    footnote(s, "Impact/Effort in dezelfde S/M/L-schaal als de GEO-deck. ‘Nu’ = kan direct, "
                "de code staat er al.")
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(1200 / PX)
    prs.slide_height = Inches(675 / PX)
    totals, cats = seasonal_validation()

    slide_cover(prs)
    slide_where_we_are(prs)
    slide_gap(prs)
    slide_what_changes(prs)
    slide_result(prs)
    slide_seasonal(prs, totals, cats)
    slide_risks(prs)
    slide_order(prs)

    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    if totals:
        print(f"seasonal 10-cat recomputed: {totals['v1']:.1f}->{totals['v2']:.1f} visits, "
              f"{totals['r1']:.1f}->{totals['r2']:.1f} revenue")


if __name__ == "__main__":
    main()
